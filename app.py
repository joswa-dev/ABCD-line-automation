import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import io
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# ---------------------------------------------------------
# 1. PAGE CONFIG & STYLING
# ---------------------------------------------------------
st.set_page_config(
    page_title="14-Makino Production & Strategy Hub",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
    <style>
    .main { background-color: #0e1117; }
    h1, h2, h3 { color: #0D5C75; font-family: 'Segoe UI', sans-serif; }
    </style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------
# TOP NAVIGATION SWITCHER
# ---------------------------------------------------------
nav_choice = st.radio(
    "Select Portal",
    ["🏭 Live 14-Makino Production Hub", "📚 Executive Presentation Gallery (20 Decks)"],
    horizontal=True
)

st.markdown("---")

# =========================================================
# MODE 1: LIVE 14-MAKINO PRODUCTION DASHBOARD
# =========================================================
if nav_choice == "🏭 Live 14-Makino Production Hub":
    st.markdown("<h2 style='text-align: center; color: #0D5C75;'>🏭 UNIT III - 14-MAKINO PRODUCTION HUB</h2>", unsafe_allow_html=True)
    st.markdown("---")

    makino_fleet = {
        # SINGLE PALLET MACHINES (1, 2, 13, 14)
        "Makino-01": {"type": "Single Pallet", "p1": "Oil Pump Housing", "p2": None, "base_p1": 210, "base_p2": 0, "dt": 15, "sc": 2, "fm": 850, "dr": 320},
        "Makino-02": {"type": "Single Pallet", "p1": "Water Pump Body", "p2": None, "base_p1": 195, "base_p2": 0, "dt": 25, "sc": 3, "fm": 910, "dr": 400},
       
        # DUAL PALLET MACHINES (3 to 12)
        "Makino-03": {"type": "Dual Pallet", "p1": "Pump Body (PB 444)", "p2": "Pressure Cover (PC)", "base_p1": 180, "base_p2": 90, "dt": 30, "sc": 4, "fm": 820, "dr": 310},
        "Makino-04": {"type": "Dual Pallet", "p1": "Transmission Case A", "p2": "Transmission Cap B", "base_p1": 160, "base_p2": 160, "dt": 20, "sc": 1, "fm": 650, "dr": 280},
        "Makino-05": {"type": "Dual Pallet", "p1": "Fuel Line Flange", "p2": "Fuel Line Valve Seat", "base_p1": 220, "base_p2": 110, "dt": 10, "sc": 2, "fm": 740, "dr": 350},
        "Makino-06": {"type": "Dual Pallet", "p1": "Manifold Intake A", "p2": "Manifold Cover B", "base_p1": 140, "base_p2": 140, "dt": 40, "sc": 5, "fm": 1050, "dr": 480},
        "Makino-07": {"type": "Dual Pallet", "p1": "Engine Mount L", "p2": "Engine Mount R", "base_p1": 175, "base_p2": 175, "dt": 15, "sc": 2, "fm": 530, "dr": 210},
        "Makino-08": {"type": "Dual Pallet", "p1": "Coolant Elbow A", "p2": "Coolant Base B", "base_p1": 200, "base_p2": 100, "dt": 0, "sc": 1, "fm": 400, "dr": 150},
        "Makino-09": {"type": "Dual Pallet", "p1": "Gearbox Side Cover", "p2": "Gearbox Plug", "base_p1": 155, "base_p2": 155, "dt": 35, "sc": 3, "fm": 980, "dr": 420},
        "Makino-10": {"type": "Dual Pallet", "p1": "Thermostat Housing", "p2": "Thermostat Cap", "base_p1": 190, "base_p2": 95, "dt": 20, "sc": 2, "fm": 790, "dr": 330},
        "Makino-11": {"type": "Dual Pallet", "p1": "Brake Bracket L", "p2": "Brake Bracket R", "base_p1": 130, "base_p2": 130, "dt": 50, "sc": 6, "fm": 1120, "dr": 470},
        "Makino-12": {"type": "Dual Pallet", "p1": "Filter Head A", "p2": "Filter Adapter B", "base_p1": 170, "base_p2": 85, "dt": 15, "sc": 1, "fm": 600, "dr": 260},
       
        # SINGLE PALLET MACHINES (13, 14)
        "Makino-13": {"type": "Single Pallet", "p1": "Heavy Flywheel Ring", "p2": None, "base_p1": 280, "base_p2": 0, "dt": 10, "sc": 1, "fm": 450, "dr": 190},
        "Makino-14": {"type": "Single Pallet", "p1": "Main Standoff Casting", "p2": None, "base_p1": 240, "base_p2": 0, "dt": 5, "sc": 0, "fm": 320, "dr": 140},
    }

    st.sidebar.header("🕹️ Shift Micro-Logging")
    selected_machine = st.sidebar.selectbox("Select Machine", list(makino_fleet.keys()))
    cfg = makino_fleet[selected_machine]

    st.sidebar.markdown(f"**Architecture:** `{cfg['type']}`")
    st.sidebar.markdown("---")

    p1_count = st.sidebar.number_input(f"{cfg['p1']} Count", min_value=0, value=cfg["base_p1"], key=f"p1_{selected_machine}")
    if cfg["type"] == "Dual Pallet":
        p2_count = st.sidebar.number_input(f"{cfg['p2']} Count", min_value=0, value=cfg["base_p2"], key=f"p2_{selected_machine}")
    else:
        p2_count = 0

    downtime_mins = st.sidebar.number_input("Unplanned Downtime (Mins)", min_value=0, value=cfg["dt"], key=f"dt_{selected_machine}")
    rejections = st.sidebar.number_input("Rejections / Scrap", min_value=0, value=cfg["sc"], key=f"sc_{selected_machine}")

    shift_length_mins = 480
    operating_time = max(0, shift_length_mins - downtime_mins)
    availability = (operating_time / shift_length_mins) * 100 if shift_length_mins > 0 else 0

    total_parts = p1_count + p2_count
    ideal_cycle_time_mins = 1.2 if cfg["type"] == "Dual Pallet" else 0.9
    performance = min(100.0, ((total_parts * ideal_cycle_time_mins) / operating_time) * 100) if operating_time > 0 else 0
    quality = ((total_parts - rejections) / total_parts * 100) if total_parts > 0 else 100
    oee = (availability / 100) * (performance / 100) * (quality / 100) * 100

    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Overall OEE", f"{oee:.1f}%", delta=f"{oee-75:.1f}% vs Target")
    col2.metric("Availability", f"{availability:.1f}%")
    col3.metric("Performance", f"{performance:.1f}%")
    col4.metric("Quality Rate", f"{quality:.1f}%")

    st.markdown("---")
    st.subheader(f"⏱️ {selected_machine} ({cfg['type']}) Shift Timeline")

    dt_end_hour = 10 + (downtime_mins // 60)
    dt_end_min = downtime_mins % 60

    timeline_data = [
        dict(Task="Production", Start="2026-07-27 06:00", Finish="2026-07-27 10:00", Status="Running"),
        dict(Task="Downtime / Setup", Start="2026-07-27 10:00", Finish=f"2026-07-27 {dt_end_hour:02d}:{dt_end_min:02d}", Status="Downtime"),
        dict(Task="Production", Start=f"2026-07-27 {dt_end_hour:02d}:{dt_end_min:02d}", Finish="2026-07-27 14:00", Status="Running"),
    ]
    df_timeline = pd.DataFrame(timeline_data)
    fig_timeline = px.timeline(
        df_timeline, x_start="Start", x_end="Finish", y="Task", color="Status",
        color_discrete_map={"Running": "#10B981", "Downtime": "#EF4444"}
    )
    fig_timeline.update_yaxes(autorange="reversed")
    fig_timeline.update_layout(height=200, margin=dict(l=10, r=10, t=10, b=10))
    st.plotly_chart(fig_timeline, use_container_width=True, config={'displayModeBar': False})

    st.markdown("---")
    col_left, col_right = st.columns(2)

    with col_left:
        st.subheader(f"📊 {selected_machine} Volume Breakdown")
        if cfg["type"] == "Dual Pallet":
            fig_bar = go.Figure(data=[
                go.Bar(name=cfg['p1'], x=['Shift Output'], y=[p1_count], text=[f"{cfg['p1']}: {p1_count}"], textposition='auto', marker_color='#0D5C75'),
                go.Bar(name=cfg['p2'], x=['Shift Output'], y=[p2_count], text=[f"{cfg['p2']}: {p2_count}"], textposition='auto', marker_color='#10B981')
            ])
            fig_bar.update_layout(barmode='stack', height=280, margin=dict(l=10, r=10, t=30, b=10))
        else:
            fig_bar = go.Figure(data=[
                go.Bar(name=cfg['p1'], x=['Shift Output'], y=[p1_count], text=[f"{cfg['p1']}: {p1_count}"], textposition='auto', marker_color='#10B981')
            ])
            fig_bar.update_layout(height=280, margin=dict(l=10, r=10, t=30, b=10))
           
        fig_bar.update_layout(
            font=dict(color="#FFFFFF"),
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            xaxis=dict(tickfont=dict(color='#FFFFFF')),
            yaxis=dict(tickfont=dict(color='#FFFFFF'))
        )
        st.plotly_chart(fig_bar, use_container_width=True, config={'displayModeBar': False})

    with col_right:
        st.subheader(f"🛠️ {selected_machine} Tool Life")
        face_mill_used = cfg["fm"] + total_parts
        driller_used = cfg["dr"] + total_parts
       
        st.write("**T01 - Main Spindle Cutter (Max 1200 Cycles)**")
        st.progress(min(1.0, face_mill_used / 1200))
        st.caption(f"Usage: {face_mill_used} / 1200 cycles ({max(0, 1200 - face_mill_used)} remaining)")
       
        st.write("**T02 - Carbide Driller (Max 500 Cycles)**")
        st.progress(min(1.0, driller_used / 500))
        st.caption(f"Usage: {driller_used} / 500 cycles ({max(0, 500 - driller_used)} remaining)")

    st.markdown("---")
    st.subheader("💾 Instant Enterprise Excel Exporter")

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
        summary_df = pd.DataFrame([{
            "Machine": selected_machine,
            "Type": cfg["type"],
            "Material 1": cfg["p1"],
            "Material 1 Count": p1_count,
            "Material 2": cfg["p2"] if cfg["p2"] else "N/A (Single Pallet)",
            "Material 2 Count": p2_count,
            "Total Output": total_parts,
            "Scrap": rejections,
            "Downtime (mins)": downtime_mins,
            "OEE (%)": round(oee, 2),
            "Availability (%)": round(availability, 2),
            "Performance (%)": round(performance, 2),
            "Quality (%)": round(quality, 2)
        }])
        summary_df.to_excel(writer, sheet_name='Makino_Report', index=False)

    excel_data = buffer.getvalue()
    st.download_button(
        label="📄 Download Production & OEE Report (.xlsx)",
        data=excel_data,
        file_name=f"Production_{selected_machine}_Report.xlsx",
        mime="application/vnd.ms-excel"
    )

# =========================================================
# MODE 2: EXECUTIVE PPT GALLERY (FULLY DYNAMIC TOPIC GENERATOR)
# =========================================================
else:
    st.markdown("<h2 style='text-align: center; color: #0D5C75;'>📚 EXECUTIVE PRESENTATION & STRATEGY GALLERY</h2>", unsafe_allow_html=True)
    st.write("Browse through executive strategy decks, line automation studies, and OEE optimization reports:")
    st.markdown("---")

    # ---------------------------------------------------------
    # DYNAMIC TOPIC-SPECIFIC SLIDE DATABASE
    # ---------------------------------------------------------
    deck_database = {
        "PPT 01: 14-Makino Line Efficiency & OEE Master Plan": {
            "summary": "Comprehensive breakdown of 14 Makino CNC machines, comparing Single vs Dual Pallet performance.",
            "slides": [
                ("14-Makino Line Efficiency & OEE Master Plan", "Unit III - Line Performance & Strategy Overview"),
                ("OEE Metric Objectives", "Targeting World-Class Efficiency Across 14 CNC Stations", "• Target Availability: 92% across all 3 operating shifts.\n• Target Performance: Maintain cycle timing under 1.2 mins.\n• Quality Rate Goal: Keep rejection scrap under 1.5%."),
                ("Line Layout Division", "Single Pallet vs Dual Pallet Architecture", "• Stations 01, 02, 13, 14: Dedicated Single-Pallet high-volume lines.\n• Stations 03 through 12: Dual-Pallet APC (Automatic Pallet Changer) setup.\n• Simultaneous cutting and part loading reduces idle times."),
                ("Availability & Downtime Breakdown", "Root Causes of Unplanned Downtimes", "• Chip clearance and spindle cleaning: 35% of minor stops.\n• Fixture loading delays & operator waiting times: 25%.\n• Unscheduled tool changes: 20%."),
                ("OEE Calculation Engine", "Mathematical Formulas Used", "• Availability = Operating Time / Planned Production Time.\n• Performance = (Total Cycles * Ideal Cycle Time) / Operating Time.\n• Quality = Passed Castings / Total Processed Castings."),
                ("Machine Heatmap & Bottlenecks", "Identifying Slow Stations", "• Makino-06 & Makino-11 show higher breakdown hours due to complex manifolds.\n• Makino-01 & 14 maintain highest uptime (98%+ Availability)."),
                ("Target vs Actual Volume Variance", "Output Benchmarking", "• Planned Weekly Output: 25,000 castings.\n• Actual Output with Micro-Logging: 23,800 castings (95.2% target rate)."),
                ("Financial Impact of OEE Losses", "Calculating Monetary Downtime Cost", "• 1 hour of unrecorded spindle breakdown = $350 lost throughput.\n• Catching 45 mins of downtime daily saves $4,500+ monthly."),
                ("Operational Improvement Plan", "Action Steps for Line Supervisors", "• Implement 5-minute shift handoff micro-logging.\n• Monitor live performance gauges twice per shift."),
                ("Executive Recommendation", "Next Steps for Senior Management", "• Standardize OEE tracking across Unit III.\n• Review weekly automated Excel exports during Monday production meetings.")
            ]
        },
        "PPT 02: Eliminating Paper Logs via Web Micro-Logging": {
            "summary": "Transitioning shift-end logging from manual paper entries to instant cloud data entry.",
            "slides": [
                ("Eliminating Paper Logs via Web Micro-Logging", "Unit III - Shift Data Automation Project"),
                ("The Paper Log Problem", "Challenges with Manual Floor Record Sheets", "• Operators spend 10-15 minutes writing notes at shift end.\n• Paper logs get damaged, stained with cutting fluid, or lost.\n• Data entry staff manually re-type numbers into Excel, introducing errors."),
                ("Web Micro-Logging Solution", "Instant Digital Entry Interface", "• Responsive web form accessible on floor tablets or smartphones.\n• Simple dropdowns and number steppers reduce input time to < 1 minute.\n• Instant validation prevents negative part counts or impossible inputs."),
                ("Database Synchronization", "Cloud vs Local Data Pipeline", "• Replaces physical binders with centralized database storage.\n• Eliminates manual data transfer delays between shop floor and office.\n• Historical logs available instantly for quality audits."),
                ("Operator Experience (UX)", "Designed for Fast Floor Input", "• Large button controls for easy tapping with work gloves.\n• High-contrast dark theme optimized for factory lighting conditions."),
                ("Error Elimination Metrics", "Before vs After Digitalization", "• Math errors on OEE calculation reduced from 14% to 0%.\n• Log completion compliance increased to 99.8%.\n• Time saved per shift: 12 minutes per operator."),
                ("IT Infrastructure & Security", "Zero-Hardware Zero-Software Footprint", "• Runs directly in any web browser without local installation.\n• Bypasses corporate IT lockdown constraints on restricted laptops."),
                ("Data Integrity & Audit Trail", "Ensuring Compliance", "• Timestamped entries prevent back-dating or false reporting.\n• Automatic calculation of scrap rates removes human bias."),
                ("Cost Comparison", "Paper/Printing Costs vs Web Portal", "• Annual paper log printing and filing cost: $2,400.\n• Cloud web dashboard operational cost: $0 (open web architecture)."),
                ("Rollout Strategy", "Transitioning All 14 Stations", "• Week 1: Dual logging (Paper + Web) on Makino-01 and 02.\n• Week 2: Full paperless shift transitions across all 14 Makinos.")
            ]
        },
        "PPT 03: Dual-Pallet APC Cycle Math & Throughput Analysis": {
            "summary": "How rotary table indexing doubles output without spindle downtime.",
            "slides": [
                ("Dual-Pallet APC Cycle Math & Throughput Analysis", "Unit III - Rotary Indexing Efficiency"),
                ("APC (Automatic Pallet Changer) Mechanics", "How Rotary Tables Work", "• Pallet A is inside the cutting chamber while Pallet B is outside.\n• Operator unloads finished part and loads raw casting on Pallet B while spindle runs.\n• 180-degree table rotation indexes Pallet B into the cutting area in < 6 seconds."),
                ("Cycle Math Comparison", "Single-Pallet vs Dual-Pallet Uptime", "• Single Pallet: Spindle stops for 45s during part reload (Spindle Uptime: 55%).\n• Dual Pallet: Spindle stops for 6s table index (Spindle Uptime: 92%)."),
                ("Throughput Calculation", "Hourly Part Production Increase", "• Single Pallet Output: 24 parts per hour.\n• Dual Pallet Output: 42 parts per hour.\n• Effective throughput increase: +75% per machine station."),
                ("Dual Material Machining Strategy", "Handling PB 444 and PC Castings", "• Pallet A cuts Pump Body (PB 444).\n• Pallet B cuts Pressure Cover (PC).\n• Dashboard tracks both parts separately in a stacked volume chart."),
                ("Indexing Wear & Preventative Maintenance", "APC Mechanism Health", "• Rotary table hydraulic pressure monitored to prevent index stalling.\n• Cumulative index cycles tracked to schedule seal replacements."),
                ("Operator Ergonomics at APC Station", "Optimizing Manual Handling", "• Standardized fixture clamping positions reduce operator wrist fatigue.\n• Pneumatic quick-clamping reduces reload times by 4 seconds."),
                ("Volume Breakdown Visualization", "Stacked Bar Chart Analytics", "• Live web view shows exact part split ratio between Pallet A and B.\n• Instantly alerts supervisor if one pallet fixture is offline."),
                ("Financial Throughput Value", "Revenue per Machine Hour", "• Dual-pallet efficiency adds 144 additional finished castings per shift.\n• Estimated added productivity value: $1,200 per station daily."),
                ("Conclusion & Best Practices", "Standard Operating Procedures", "• Mandate simultaneous loading during active spindle cuts.\n• Maintain clean rotary table surface to prevent indexing alignment faults.")
            ]
        },
        "PPT 04: Predictive Tool Life & Spindle Wear Modeling": {
            "summary": "Tracking cutter impact cycles mathematically to schedule tool changes before failure.",
            "slides": [
                ("Predictive Tool Life & Spindle Wear Modeling", "Unit III - Tool Wear & Spindle Maintenance"),
                ("The Risk of Unexpected Tool Failure", "Why Tool Snap Occurs", "• Cutting tools wear down progressively with every impact cycle on aluminum castings.\n• Worn tools cause dimensional errors, poor surface finish, or sudden snapping.\n• Snapped tool inside casting damages spindle and causes 2+ hours downtime."),
                ("Cycle Counter Mathematical Model", "Tracking Wear Without Sensors", "• Every entered part count deducts from the tool's rated cycle lifespan.\n• T01 Main Spindle Face Mill: Maximum 1,200 cycles.\n• T02 Carbide Drill Bit: Maximum 500 cycles."),
                ("Progress Bar Warning System", "Visual Alerts for Floor Team", "• Green Bar: 0% - 70% Tool Life used (Safe Operation).\n• Yellow Warning: 71% - 90% Tool Life used (Prepare Replacement Tool).\n• Red Alert: 90%+ Tool Life used (Mandatory Change Before Next Shift)."),
                ("Spindle Load & Overhaul Scheduling", "Tracking Running Hours", "• Cumulative cutting hours tracked mathematically to predict spindle bearing wear.\n• Major spindle overhaul scheduled automatically at 4,000 running hours."),
                ("Scrap Reduction via Tool Health", "Quality Correlation", "• 80% of casting dimensional rejections occur in the final 10% of tool life.\n• Changing tools proactively at 95% cycle limit reduces scrap by 68%."),
                ("Cost Analysis: Proactive vs Reactive Tooling", "Financial Savings", "• Reactive (Tool breaks mid-cut): Tool cost ($120) + Scrap ($45) + Spindle repair ($800) = $965.\n• Proactive (Changed at cycle limit): Tool cost ($120) = $120."),
                ("Operator Tool Change Workflow", "Simple Reset Button", "• Operator replaces worn tool, clicks 'Reset Counter' on dashboard.\n• System logs timestamp, tool ID, and operator ID for tracking."),
                ("Integration with Excel Exporters", "Maintenance Reporting", "• Daily Excel exports include remaining tool life counts for maintenance team."),
                ("Executive Summary", "Tool Maintenance Mandate", "• Mandate cycle-based tool replacements across all 14 Makino stations.")
            ]
        }
    }

    # Fallback template generator for decks 05 through 20
    for i in range(5, 21):
        deck_key = f"PPT {i:02d}: Executive Strategy Deck {i:02d}"
        if i == 5:
            deck_key = "PPT 05: Non-IoT Digital Twin Architecture"
            summ = "Achieving $100k industrial visibility without expensive PLC sensors or hardware modifications."
            topic = "Non-IoT Digital Twin Architecture"
        elif i == 6:
            deck_key = "PPT 06: Unplanned Downtime Root Cause Analysis"
            summ = "Category breakdown of chip clearance delays, maintenance stops, and power cuts."
            topic = "Unplanned Downtime Root Cause Analysis"
        elif i == 7:
            deck_key = "PPT 07: Quality & Rejection Rate Optimization"
            summ = "Scrap reduction techniques across Pump Body and Pressure Cover castings."
            topic = "Quality & Rejection Rate Optimization"
        elif i == 8:
            deck_key = "PPT 08: Shop Floor Ergonomics & Shift Transition"
            summ = "Reducing 10-minute shift end paperwork burden for shop operators."
            topic = "Shop Floor Ergonomics & Shift Transition"
        elif i == 9:
            deck_key = "PPT 09: Automated Shift Report Generation"
            summ = "Creating structured Excel audit logs directly from micro-log database entries."
            topic = "Automated Shift Report Generation"
        elif i == 10:
            deck_key = "PPT 10: Executive Summary - Line Automation ROI"
            summ = "Cost-benefit analysis of web dashboard vs traditional manual entry."
            topic = "Line Automation ROI Analysis"
        elif i == 11:
            deck_key = "PPT 11: Single-Pallet vs Dual-Pallet Machine Layout"
            summ = "Operational comparison of Makino 01, 02, 13, 14 vs Makino 03-12."
            topic = "Single vs Dual Pallet Layout Comparison"
        elif i == 12:
            deck_key = "PPT 12: Production Target vs Actual Volume Variance"
            summ = "Weekly metrics and gap analysis for line supervisors."
            topic = "Production Target vs Volume Variance"
        elif i == 13:
            deck_key = "PPT 13: Spindle Speed & Feed Rate Optimization"
            summ = "Cutting parameter adjustments for high-speed aluminum casting milling."
            topic = "Spindle Speed & Feed Rate Optimization"
        elif i == 14:
            deck_key = "PPT 14: Preventive Maintenance Scheduling Framework"
            summ = "Transitioning from reactive tool replacement to proactive cycle-based maintenance."
            topic = "Preventive Maintenance Framework"
        elif i == 15:
            deck_key = "PPT 15: Shift A, B, C Comparative Line Analytics"
            summ = "Cross-shift efficiency and performance distribution models."
            topic = "Shift A, B, C Comparative Analytics"
        elif i == 16:
            deck_key = "PPT 16: Zero-Cost Factory Digitalization Strategy"
            summ = "How web technologies bypass strict corporate IT laptop lockdown rules."
            topic = "Zero-Cost Factory Digitalization"
        elif i == 17:
            deck_key = "PPT 17: Machine Line Heatmap & Bottleneck Identification"
            summ = "Pinpointing slow cycles across all 14 CNC stations."
            topic = "Line Heatmap & Bottleneck Identification"
        elif i == 18:
            deck_key = "PPT 18: Scrap Reclamation & Material Handling"
            summ = "Tracking raw casting defects vs machining scrap."
            topic = "Scrap Reclamation & Material Handling"
        elif i == 19:
            deck_key = "PPT 19: Operator User Experience & Mobile Web Integration"
            summ = "Optimizing dashboard UI for Android shop-floor tablets."
            topic = "Operator UX & Mobile Integration"
        else:
            deck_key = "PPT 20: Future Roadmap - Non-Invasive Optical Sensors"
            summ = "Next-level external door sensor integration for fully automated cycle counts."
            topic = "Future Roadmap & Optical Sensors"

        deck_database[deck_key] = {
            "summary": summ,
            "slides": [
                (f"{topic} Overview", f"Unit III - Executive Strategy Deck {i:02d}"),
                ("Primary Objectives", f"Key goals for {topic}", f"• Target 1: Improve line visibility across all 14 stations.\n• Target 2: Reduce operational delays and manual friction.\n• Target 3: Establish standardized metrics for shift supervisors."),
                ("Current Baseline Analysis", "Benchmarking Current Floor Performance", f"• Historical data analysis for {topic}.\n• Identifying primary operational bottlenecks.\n• Operator feedback and shop floor observations."),
                ("Proposed Solution Framework", "Core Methodologies Implemented", f"• Integrating {topic} into daily line routines.\n• Automated data processing via Python analytics engine.\n• High-contrast visual dashboards for real-time tracking."),
                ("Performance Impact & Metrics", "Expected Efficiency Gains", f"• Projected throughput increase: +12% to +18%.\n• Reduced calculation errors from 10% to 0%.\n• Instant compliance reporting for management audits."),
                ("Implementation Workflow", "Step-by-Step Operator Guidelines", "• Step 1: Access web portal on station tablet.\n• Step 2: Input micro-log data at shift end.\n• Step 3: Review dynamic OEE calculations instantly."),
                ("Cost & Resource Analysis", "Financial Optimization Model", f"• Zero additional hardware required.\n• Utilizing open web architecture and cloud deployment.\n• ROI realized within first month of deployment."),
                ("Risk Mitigation", "Addressing Operational Challenges", "• Backup offline logging capabilities.\n• Automated input validation to prevent data corruption.\n• Simple UI requiring zero technical training."),
                ("Rollout Schedule", "Deployment Roadmap", "• Phase 1: Pilot testing on select Makino stations.\n• Phase 2: Full line integration across all 14 stations.\n• Phase 3: Weekly management review and continuous tuning."),
                ("Executive Action Items", "Next Steps for Senior Leadership", f"• Review and approve {topic} implementation.\n• Authorize digital micro-logging protocol.\n• Schedule weekly OEE audit reviews.")
            ]
        }

    selected_deck = st.selectbox("Select Presentation Deck", list(deck_database.keys()))
    deck_info = deck_database[selected_deck]
   
    st.info(f"**Deck Summary:** {deck_info['summary']}")

    col_view, col_info = st.columns([2, 1])

    with col_view:
        st.subheader("🖼️ Interactive Deck Preview (10 Topic-Specific Slides)")
       
        slides_preview_html = "<br/>".join([f"• <b>Slide {idx+1}:</b> {slide[0]}" for idx, slide in enumerate(deck_info["slides"])])
       
        st.markdown(f"""
        <div style="background-color: #1e293b; padding: 25px; border-radius: 10px; border: 1px solid #334155; color: white;">
            <h3 style="color: #38bdf8;">📊 {selected_deck}</h3>
            <p style="color: #94a3b8;"><b>Total Presentation Deck Length:</b> 10 Custom Slides</p>
            <hr style="border-color: #334155;"/>
            <div style="background-color: #0f172a; padding: 15px; border-radius: 6px; font-family: sans-serif; font-size: 14px; line-height: 1.8;">
                <b>Slide Index Included in Download:</b><br/>
                {slides_preview_html}
            </div>
        </div>
        """, unsafe_allow_html=True)

    with col_info:
        st.subheader("💾 File Actions")
        st.write("Download the customized 10-slide PowerPoint file:")
       
        # ---------------------------------------------------------
        # GENERATE CUSTOM TOPIC PPTX BINARY
        # ---------------------------------------------------------
        prs = Presentation()
        bullet_layout = prs.slide_layouts[1]

        for idx, item in enumerate(deck_info["slides"]):
            slide = prs.slides.add_slide(bullet_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
           
            if len(item) == 2:
                # Title slide layout
                header, sub_header = item
                content = ""
            else:
                header, sub_header, content = item

            title_shape.text = f"Slide {idx+1}: {header}"
            tf = body_shape.text_frame
            tf.text = sub_header
            if content:
                p = tf.add_paragraph()
                p.text = content

        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        clean_filename = selected_deck.split(':')[0].replace(' ', '_') + "_Custom10Slides.pptx"

        st.download_button(
            label=f"📥 Download Customized 10-Slide Deck (.pptx)",
            data=ppt_buffer.getvalue(),
            file_name=clean_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
       
        st.markdown("---")
        st.markdown("### 📌 Deck Specs:")
        st.write("• **Total Slides:** 10 Topic-Specific Slides")
        st.write("• **Format:** Genuine Microsoft PowerPoint (.pptx)")
        st.write("• **Target:** Management & Senior Engineers") 
